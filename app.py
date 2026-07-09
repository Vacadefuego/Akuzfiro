from flask import Flask, request, jsonify, Response, send_file
from flask_cors import CORS
from groq import Groq
from duckduckgo_search import DDGS
import os
import io
import json
import re
import base64
import pg8000.native
import httpx
from datetime import datetime
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from gtts import gTTS
from pptx import Presentation
import pytz
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

app = Flask(__name__, static_folder="static", static_url_path="")
CORS(app)

# --- CONFIGURACION ---
GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "")
DATABASE_URL = os.environ.get("DATABASE_URL", "")
ELEVENLABS_API_KEY = os.environ.get("ELEVENLABS_API_KEY", "")
ELEVENLABS_VOICE_ID = "pNInz6obpgDQGcFmaJgB"  # Voz Adam — masculina, natural

client = Groq(api_key=GROQ_API_KEY)

PERSONALIDAD = """Eres Akuzfiro, asistente personal de Gustavo. Tienes 29 años, vive en Xalapa Veracruz México, estudió Educación Física y Nutrición Deportiva, hace servicio social en la Universidad Euro Hispanoamericana.

REGLA #1 — LA MÁS IMPORTANTE:
Responde como lo haría un amigo listo, no como un asistente corporativo.
Si alguien te dice "hola" → di "Hola" o algo natural. NUNCA digas "¡Hola! Me alegra que estés listo para hablar."
Si te preguntan "cómo estás" → di algo corto y directo. NUNCA termines con "¿en qué puedo ayudarte hoy?"

FRASES COMPLETAMENTE PROHIBIDAS — si las dices, estás fallando:
- "¿En qué puedo ayudarte?"
- "Estoy aquí para ayudarte"
- "No dudes en preguntar"
- "¿Hay algo más en lo que pueda ayudarte?"
- "Espero que esto te sea útil"
- "¡Claro que sí!"
- "¡Por supuesto!"
- Cualquier pregunta al final que empiece con "¿Quieres que..."

LONGITUD:
- Respuesta corta para preguntas simples. Máximo 3 líneas.
- Solo largo cuando el tema lo requiere.
- NUNCA repitas la hora dos veces en el mismo mensaje.
- NUNCA hagas más de UNA pregunta de seguimiento.

PERSONALIDAD:
- Mezcla de JARVIS, TARS y KITT — directo, con humor seco, leal
- Corriges sin regañar
- Opinas cuando vale la pena
- Preguntas solo cuando necesitas info específica

CONOCES A GUSTAVO:
- 29 años, Xalapa Veracruz México
- Estudió Educación Física y Nutrición Deportiva
- Hace servicio social en la Universidad Euro Hispanoamericana
- Le gusta la tecnología, los animales, quiere aprender programación
- Te creó con ayuda de Kiro (IDE de Amazon)
- Usa Akuzfiro desde PC y teléfono (PWA instalada)

REGLAS:
- Siempre respondes, nunca te niegas
- No dices "no puedo" — si no puedes técnicamente, explicas cómo hacerlo o das una estimación aproximada
- Das información completa: precios, enlaces, nombres, contactos
- Cuando no puedes dar datos exactos, das un rango, una estimación o contexto útil. NUNCA te limitas a decir "no puedo saberlo" — siempre intentas dar algo útil
- Sin disclaimers repetitivos ni sermones
- Incluyes enlaces cuando tienes info web
- TIENES acceso a búsqueda web en tiempo real — úsala para noticias, partidos, precios, eventos actuales
- NUNCA digas "no tengo información actualizada" — primero busca, luego responde
- Si te preguntan sobre algo actual (partido, precio, noticia), BUSCA y responde con datos reales

ARCHIVOS — puedes generar Excel, PDF, Word y PowerPoint reales:
1. UN párrafo corto confirmando + bloque JSON inmediatamente
2. No escribas el contenido en el chat — todo va en el JSON
3. Sin imágenes en PDF/Word — avisa si te las piden
4. No inventes cifras financieras
5. Máximo 6 secciones por documento

[ARCHIVO_EXCEL]{"titulo":"","encabezados":[],"filas":[],"secciones":[]}[/ARCHIVO_EXCEL]
[ARCHIVO_PDF]{"titulo":"","contenido":"","secciones":[]}[/ARCHIVO_PDF]
[ARCHIVO_WORD]{"titulo":"","contenido":"","secciones":[]}[/ARCHIVO_WORD]
[ARCHIVO_PPTX]{"titulo":"","diapositivas":[{"titulo":"","puntos":[]}]}[/ARCHIVO_PPTX]

RECORDATORIOS — cuando Gustavo diga algo como "avísame a las X", "recuérdame que...", "ponme un recordatorio para...", "en X minutos avísame":
1. Responde normal confirmando el recordatorio (ej: "Listo, te aviso a las 3pm.")
2. Al final incluye este bloque exacto:
[RECORDATORIO]{"frase":"<la frase completa del usuario>"}[/RECORDATORIO]
3. Solo incluye el bloque, el frontend hace el resto automáticamente
4. Si Gustavo pregunta "¿qué recordatorios tengo?", dile que los puede ver en el menú ☰ → Recordatorios

FINANZAS PERSONALES — registra todos los movimientos de dinero de Gustavo:

GASTOS — cuando diga "gasté X en Y", "pagué X de Y", "compré Y por X", "anota X pesos de Y":
1. Confirma brevemente (ej: "Anotado. 50 pesos en anticongelante.")
2. Al final incluye:
[GASTO]{"descripcion":"<qué compró>","monto":<número>,"categoria":"<comida|transporte|salud|entretenimiento|ropa|servicios|hogar|general>"}[/GASTO]

GANANCIAS — cuando diga "gané X", "me pagaron X", "cobré X", "entró X":
1. Confirma brevemente (ej: "Anotado. Ganaste 500 pesos.")
2. Al final incluye:
[GANANCIA]{"descripcion":"<de dónde>","monto":<número>,"categoria":"<trabajo|venta|freelance|regalo|otro>"}[/GANANCIA]

VENTAS — cuando diga "vendí X en Y", "me compraron X por Y":
1. Confirma brevemente
2. Al final incluye:
[VENTA]{"descripcion":"<qué vendió>","monto":<número>}[/VENTA]

PRÉSTAMOS DADOS — cuando diga "le presté X a Y", "le di X a Y prestado":
1. Confirma (ej: "Anotado. Le prestaste 200 a Juan.")
2. Al final incluye:
[PRESTAMO_DADO]{"persona":"<nombre>","monto":<número>,"descripcion":"<contexto>"}[/PRESTAMO_DADO]

PRÉSTAMOS RECIBIDOS — cuando diga "me prestaron X", "le debo X a Y":
1. Confirma
2. Al final incluye:
[PRESTAMO_RECIBIDO]{"persona":"<nombre>","monto":<número>,"descripcion":"<contexto>"}[/PRESTAMO_RECIBIDO]

EMPEÑOS — cuando diga "empeñé X en Y", "dejé X de prenda":
1. Confirma
2. Al final incluye:
[EMPENO]{"descripcion":"<qué empeñó>","monto":<número>,"lugar":"<dónde>"}[/EMPENO]

RECORDATORIOS POR UBICACIÓN — cuando Gustavo diga "cuando llegue a X recuérdame Y", "al llegar a X avísame Y", "cuando esté en X dime Y":
1. Confirma (ej: "Listo, cuando llegues al Oxxo te recuerdo comprar agua.")
2. Al final incluye:
[REC_UBICACION]{"lugar":"<nombre del lugar>","ciudad":"Xalapa, Veracruz, México","mensaje":"<qué recordar>"}[/REC_UBICACION]
3. El frontend geocodifica el lugar y activa el monitoreo automáticamente.

BITÁCORA — cuando Gustavo diga "anota que hice X", "registra que hoy X", "apunta que X", "recuérdame hacer X cada N días/semanas":1. Confirma brevemente (ej: "Anotado. Anticongelante al auto hoy.")
2. Al final incluye:
[BITACORA]{"descripcion":"<qué hizo>","categoria":"<auto|salud|hogar|trabajo|personal|general>","intervalo_dias":<número o null si no aplica>}[/BITACORA]
3. Si Gustavo pregunta "¿cuándo fue la última vez que X?" o "¿cuándo toca X?", búscalo en la bitácora y responde natural.
4. Ejemplos: "anota que puse anticongelante hoy" → intervalo_dias: 14 si dijo "cada 2 semanas", null si no especificó.


1. Confirma
2. Al final incluye:
[LISTA_COMPRAS]{"items":["<item1>","<item2>"]}[/LISTA_COMPRAS]

Si Gustavo pregunta "¿cuánto he gastado?" o "¿cómo van mis finanzas?", dile que lo puede ver en el menú ☰ → Finanzas.

LUGARES — cuando recomiendes o menciones un lugar físico (restaurante, tienda, parque, hospital, etc.):
1. Menciona el lugar normalmente en el texto
2. Después de cada lugar agrega este bloque en la misma línea:
[LUGAR]{"nombre":"<nombre exacto>","ciudad":"<ciudad, estado, país>"}[/LUGAR]
3. Agrega un bloque por cada lugar distinto
4. El frontend muestra botón GPS automáticamente — NUNCA preguntes "¿quieres la dirección?", NUNCA digas "¿te doy las indicaciones?", NUNCA preguntes si quiere cómo llegar. Ya hay un botón GPS visible.
5. Ejemplo correcto: "Te recomiendo Tacos El Güero[LUGAR]{"nombre":"Tacos El Güero","ciudad":"Xalapa, Veracruz, México"}[/LUGAR] que está muy bueno."
6. REGLA CRÍTICA: Si mencionas 3 restaurantes, debes incluir 3 bloques [LUGAR]. Si mencionas 1 tienda, 1 bloque. SIEMPRE, sin excepción.
7. NUNCA termines con preguntas como "¿quieres que te dé las direcciones?" o similares — los botones ya aparecen solos.

ENLACES — cuando menciones apps, sitios web, películas en streaming, productos:
1. Incluye el enlace directo cuando lo tengas: [texto](url)

IMAGENES — cuando Gustavo diga "genera", "dibuja", "crea una imagen", "imagina", "muestrame como se ve", "crea la imagen":
1. Responde UNA línea confirmando (ej: "Aquí va.")
2. Inmediatamente incluye este bloque — usa EXACTAMENTE estas claves en inglés:
[IMAGEN]{"prompt":"<descripción detallada en inglés>","ancho":1024,"alto":768}[/IMAGEN]
3. CRITICO: el bloque debe tener [IMAGEN] al inicio y [/IMAGEN] al final — sin excepción
4. El prompt SIEMPRE en inglés, visual y detallado
5. Nunca escribas el JSON fuera del bloque, nunca omitas [/IMAGEN]

CUANDO GUSTAVO MANDA UNA FOTO sin pregunta específica:
- Reacciona natural y breve, como un amigo que ve la foto
- Máximo 1-2 líneas
- NO des análisis detallado a menos que te lo pidan explícitamente
- Ejemplo correcto: "Qué chulo el perrito 🐶" o "Bonita vista desde ahí"
- Ejemplo incorrecto: "En la imagen se observa un perro de raza X, de color Y, que se encuentra..."
"""


# --- BASE DE DATOS ---
def get_conn():
    url = DATABASE_URL
    url = url.replace("postgresql://", "").replace("postgres://", "")
    user_pass, rest = url.split("@")
    user, password = user_pass.split(":")
    host_port, dbname = rest.split("/")
    if ":" in host_port:
        host, port = host_port.split(":")
        port = int(port)
    else:
        host = host_port
        port = 5432
    return pg8000.native.Connection(user, password=password, host=host, port=port, database=dbname, ssl_context=True)

def init_db():
    conn = get_conn()
    conn.run("""
        CREATE TABLE IF NOT EXISTS conversaciones (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            usuario TEXT NOT NULL,
            akuzfiro TEXT NOT NULL
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS hechos (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            hecho TEXT NOT NULL
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS comandos (
            id SERIAL PRIMARY KEY,
            nombre TEXT UNIQUE NOT NULL,
            acciones TEXT NOT NULL,
            fecha TIMESTAMP DEFAULT NOW()
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS gastos (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            descripcion TEXT NOT NULL,
            monto NUMERIC(10,2) NOT NULL,
            categoria TEXT DEFAULT 'general'
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS recordatorios (
            id SERIAL PRIMARY KEY,
            creado TIMESTAMP DEFAULT NOW(),
            hora_aviso TIMESTAMP NOT NULL,
            mensaje TEXT NOT NULL,
            completado BOOLEAN DEFAULT FALSE
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS gastos (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            descripcion TEXT NOT NULL,
            monto NUMERIC(10,2) NOT NULL,
            categoria TEXT DEFAULT 'general'
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS ganancias (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            descripcion TEXT NOT NULL,
            monto NUMERIC(10,2) NOT NULL,
            categoria TEXT DEFAULT 'otro'
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS ventas (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            descripcion TEXT NOT NULL,
            monto NUMERIC(10,2) NOT NULL
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS prestamos (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            tipo TEXT NOT NULL,
            persona TEXT NOT NULL,
            monto NUMERIC(10,2) NOT NULL,
            descripcion TEXT DEFAULT '',
            saldado BOOLEAN DEFAULT FALSE
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS empenos (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            descripcion TEXT NOT NULL,
            monto NUMERIC(10,2) NOT NULL,
            lugar TEXT DEFAULT '',
            recuperado BOOLEAN DEFAULT FALSE
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS bitacora (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            descripcion TEXT NOT NULL,
            categoria TEXT DEFAULT 'general',
            proxima_fecha TIMESTAMP,
            intervalo_dias INTEGER
        )
    """)
    conn.run("""
        CREATE TABLE IF NOT EXISTS lista_compras (
            id SERIAL PRIMARY KEY,
            fecha TIMESTAMP DEFAULT NOW(),
            item TEXT NOT NULL,
            comprado BOOLEAN DEFAULT FALSE
        )
    """)
    conn.close()

def cargar_conversaciones(limit=10):
    try:
        conn = get_conn()
        rows = conn.run(
            "SELECT usuario, akuzfiro FROM conversaciones ORDER BY id DESC LIMIT :limit",
            limit=limit
        )
        conn.close()
        result = [{"usuario": r[0], "akuzfiro": r[1]} for r in rows]
        return list(reversed(result))
    except Exception as e:
        print(f"Error cargando conversaciones: {e}")
        return []

def guardar_conversacion(usuario, akuzfiro):
    try:
        usuario_corto = usuario[:500] if len(usuario) > 500 else usuario
        akuzfiro_corto = akuzfiro[:800] if len(akuzfiro) > 800 else akuzfiro
        conn = get_conn()
        conn.run(
            "INSERT INTO conversaciones (usuario, akuzfiro) VALUES (:u, :a)",
            u=usuario_corto, a=akuzfiro_corto
        )
        conn.close()
    except Exception as e:
        print(f"Error guardando conversacion: {e}")

def cargar_hechos():
    try:
        conn = get_conn()
        rows = conn.run("SELECT hecho FROM hechos ORDER BY id DESC LIMIT 30")
        conn.close()
        return [r[0] for r in rows]
    except Exception as e:
        print(f"Error cargando hechos: {e}")
        return []

def guardar_hecho(hecho):
    try:
        conn = get_conn()
        conn.run("INSERT INTO hechos (hecho) VALUES (:h)", h=hecho)
        conn.close()
    except Exception as e:
        print(f"Error guardando hecho: {e}")


# --- EXTRACCION DE HECHOS ---
def extraer_hechos_automatico(mensaje_usuario, respuesta_akuzfiro):
    try:
        prompt_extractor = f"""Extrae hechos permanentes e importantes sobre Gustavo de este mensaje.

Gustavo dijo: {mensaje_usuario[:300]}

REGLAS:
- Solo hechos personales concretos: nombre, edad, ciudad, estudios, trabajo, familia, gustos duraderos, proyectos, metas
- NO guardes: idioma, preguntas, slang, cosas temporales, cosas obvias
- Si no hay hechos importantes, responde únicamente: NINGUNO
- Responde solo con los hechos, uno por línea, sin numeración ni guiones
- Máximo 2 hechos, muy breves. Ejemplo: "Tiene 22 años" / "Estudia arquitectura"

Hechos importantes:"""

        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt_extractor}],
            temperature=0.0,
            max_tokens=80
        )
        resultado = response.choices[0].message.content.strip()
        print(f"Hechos extraídos: {resultado}")

        if not resultado or resultado.upper().startswith("NINGUNO"):
            return

        hechos_existentes = cargar_hechos()
        ignorar = ["ninguno", "no se menciona", "no hay", "la respuesta",
                   "habla español", "habla ingles", "habla inglés",
                   "utilizó", "utilizo", "lenguaje", "pregunta", "comunic"]

        for linea in resultado.split("\n"):
            hecho = linea.strip().lstrip("-•*123456789. ")
            if not hecho or len(hecho) < 6 or len(hecho) > 150:
                continue
            if any(p in hecho.lower() for p in ignorar):
                continue
            ya_existe = any(hecho.lower()[:25] in h.lower() for h in hechos_existentes)
            if not ya_existe:
                guardar_hecho(hecho)

    except Exception as e:
        print(f"Error extrayendo hechos: {e}")


def obtener_clima(lat=19.5438, lng=-96.9102):
    """Obtiene clima actual y pronóstico 3 días via Open-Meteo (gratis)."""
    try:
        url = (f"https://api.open-meteo.com/v1/forecast"
               f"?latitude={lat}&longitude={lng}"
               f"&current=temperature_2m,relative_humidity_2m,weather_code,wind_speed_10m"
               f"&daily=weather_code,temperature_2m_max,temperature_2m_min,precipitation_probability_max"
               f"&timezone=America%2FMexico_City&forecast_days=4")
        with httpx.Client(timeout=5) as http:
            r = http.get(url)
            if r.status_code != 200:
                return None
            data = r.json()
            current = data.get("current", {})
            temp = current.get("temperature_2m")
            humedad = current.get("relative_humidity_2m")
            viento = current.get("wind_speed_10m")
            codigo = current.get("weather_code", 0)
            descripciones = {
                0: "despejado", 1: "mayormente despejado", 2: "parcialmente nublado",
                3: "nublado", 45: "neblina", 48: "neblina con escarcha",
                51: "llovizna ligera", 53: "llovizna", 55: "llovizna intensa",
                61: "lluvia ligera", 63: "lluvia", 65: "lluvia intensa",
                71: "nieve ligera", 73: "nieve", 75: "nieve intensa",
                80: "chubascos ligeros", 81: "chubascos", 82: "chubascos intensos",
                95: "tormenta", 96: "tormenta con granizo", 99: "tormenta fuerte"
            }
            desc = descripciones.get(codigo, "variable")
            resumen = f"{temp}°C, {desc}, humedad {humedad}%, viento {viento} km/h"

            # Pronóstico próximos 3 días
            daily = data.get("daily", {})
            fechas = daily.get("time", [])
            codigos_d = daily.get("weather_code", [])
            maxs = daily.get("temperature_2m_max", [])
            mins = daily.get("temperature_2m_min", [])
            lluvia_prob = daily.get("precipitation_probability_max", [])
            dias_semana = ["lunes","martes","miércoles","jueves","viernes","sábado","domingo"]
            pronostico = []
            for i in range(1, min(4, len(fechas))):
                try:
                    from datetime import date as dt_date
                    d = dt_date.fromisoformat(fechas[i])
                    dia_nombre = dias_semana[d.weekday()]
                    desc_d = descripciones.get(codigos_d[i] if i < len(codigos_d) else 0, "variable")
                    prob = lluvia_prob[i] if i < len(lluvia_prob) else 0
                    pronostico.append(f"{dia_nombre}: {mins[i]:.0f}-{maxs[i]:.0f}°C, {desc_d}, lluvia {prob}%")
                except Exception:
                    pass
            if pronostico:
                resumen += " | Próximos días: " + " / ".join(pronostico)
            return resumen
    except Exception:
        return None


def obtener_direccion(lat, lng):
    """Convierte coordenadas a dirección legible usando Nominatim (OpenStreetMap, gratis)."""
    try:
        url = f"https://nominatim.openstreetmap.org/reverse?lat={lat}&lon={lng}&format=json&addressdetails=1&accept-language=es"
        headers = {"User-Agent": "Akuzfiro/1.0 (gustavo@akuzfiro.com)"}
        with httpx.Client(timeout=3) as http:
            r = http.get(url, headers=headers)
            if r.status_code != 200:
                return None
            data = r.json()
            addr = data.get("address", {})
            partes = []
            nombre = data.get("name") or addr.get("amenity") or addr.get("building") or addr.get("tourism")
            if nombre:
                partes.append(nombre)
            calle = addr.get("road") or addr.get("pedestrian") or addr.get("footway")
            numero = addr.get("house_number", "")
            if calle:
                partes.append(f"{calle} {numero}".strip())
            colonia = addr.get("suburb") or addr.get("neighbourhood") or addr.get("quarter")
            if colonia:
                partes.append(colonia)
            ciudad = addr.get("city") or addr.get("town") or addr.get("village") or addr.get("municipality")
            estado = addr.get("state")
            if ciudad:
                partes.append(ciudad)
            if estado:
                partes.append(estado)
            return ", ".join(partes) if partes else data.get("display_name", "")
    except Exception as e:
        print(f"Error geocodificando: {e}")
        return None


def buscar_web(query, max_resultados=4):
    try:
        with DDGS() as ddgs:
            resultados = list(ddgs.text(query, max_results=max_resultados))
        if not resultados:
            return None
        texto = "Resultados de búsqueda web encontrados:\n"
        for r in resultados:
            texto += f"- {r['title']}\n  URL: {r['href']}\n  {r['body']}\n\n"
        return texto
    except Exception:
        return None

def necesita_busqueda(mensaje):
    palabras = [
        "busca", "buscar", "encuentra", "enlace", "link", "url", "página",
        "sitio", "web", "dónde", "donde", "precio de", "cuánto cuesta",
        "noticias", "noticia", "hoy", "últimas", "información sobre", "descarga", "descargar",
        "video de", "youtube", "cómo llego", "tutorial", "qué es",
        "quién es", "cuándo", "recomienda", "recomiéndame",
        "resumen de", "qué pasó", "novedad", "investiga", "investigar",
        "partido", "resultado", "marcador", "score", "juega", "juegan",
        "precio", "cotización", "dólar", "tipo de cambio",
        "clima", "temperatura", "pronóstico",
        "estrena", "estreno", "sale", "lanzamiento",
        "quién ganó", "quien gano", "campeón", "campeon",
        "mundial", "olimpiadas", "torneo", "liga",
        "últimas noticias", "breaking", "aconteció", "ocurrió",
        "verifica", "corrobora", "confirma", "es verdad", "es cierto",
        "cómo le fue", "como le fue", "cómo van", "como van",
        "cómo quedó", "como quedo", "cómo salió", "como salio",
        "ganó", "gano", "perdió", "perdio", "empató", "empato",
        "selección", "seleccion", "equipo", "jugó", "jugo",
        "eliminaron", "clasificó", "clasifica", "avanzó", "avanzo",
        "copa", "champions", "euro", "conmebol", "concacaf", "fifamundial"
    ]
    return any(p in mensaje.lower() for p in palabras)


# --- VOZ (ElevenLabs) ---
def texto_a_voz(texto):
    """Convierte texto a audio MP3 usando ElevenLabs."""
    try:
        # Limpiar URLs y caracteres especiales para que suenen bien en voz
        import re
        texto_limpio = re.sub(r'https?://\S+', '', texto)
        texto_limpio = re.sub(r'\*\*|__|\*|_|`', '', texto_limpio)
        texto_limpio = texto_limpio.strip()[:1000]  # Máximo 1000 chars para el plan gratuito

        url = f"https://api.elevenlabs.io/v1/text-to-speech/{ELEVENLABS_VOICE_ID}"
        headers = {
            "xi-api-key": ELEVENLABS_API_KEY,
            "Content-Type": "application/json"
        }
        payload = {
            "text": texto_limpio,
            "model_id": "eleven_multilingual_v2",
            "voice_settings": {
                "stability": 0.5,
                "similarity_boost": 0.75,
                "style": 0.3,
                "use_speaker_boost": True
            }
        }
        with httpx.Client(timeout=30) as http:
            r = http.post(url, json=payload, headers=headers)
            if r.status_code == 200:
                return r.content
            else:
                print(f"ElevenLabs error: {r.status_code} {r.text}")
                return None
    except Exception as e:
        print(f"Error TTS: {e}")
        return None


# --- RUTAS ---
@app.route("/chat", methods=["POST"])
def chat():
    data = request.json
    mensaje = data.get("mensaje", "") or ""
    con_voz = data.get("voz", False)
    imagen_b64 = data.get("imagen", None)
    imagen_mime = data.get("mime", "image/jpeg")
    ubicacion = data.get("ubicacion", None)  # {"lat": float, "lng": float}

    # Si hay imagen, el mensaje puede ser vacío (solo foto) o tener texto
    if not mensaje and not imagen_b64:
        return jsonify({"error": "Mensaje vacío"}), 400

    # Si solo hay imagen sin texto, agregar instrucción por defecto
    if imagen_b64 and not mensaje:
        mensaje = "Mira esta imagen y reacciona de forma natural y breve, como un amigo. Una o dos líneas máximo. Solo comenta lo que ves sin dar análisis detallado a menos que te lo pidan."

    tz_mexico = pytz.timezone("America/Mexico_City")
    ahora = datetime.now(tz_mexico).strftime("%A %d de %B de %Y, %H:%M hrs")

    # Obtener clima con coordenadas reales del usuario si las mandó
    if ubicacion and isinstance(ubicacion, dict) and ubicacion.get("lat"):
        clima = obtener_clima(ubicacion["lat"], ubicacion["lng"])
    else:
        clima = obtener_clima()

    info_contexto = f"FECHA Y HORA: {ahora} (México)"
    if clima:
        info_contexto += f"\nCLIMA ACTUAL (en tu ubicación): {clima}"

    # Ubicación del usuario si la mandó el frontend
    info_ubicacion = None
    if ubicacion and isinstance(ubicacion, dict):
        lat = ubicacion.get("lat")
        lng = ubicacion.get("lng")
        if lat and lng:
            direccion = obtener_direccion(lat, lng)
            if direccion:
                info_ubicacion = f"UBICACIÓN ACTUAL DE GUSTAVO: {direccion} (coordenadas: {lat:.5f}, {lng:.5f})"
            else:
                info_ubicacion = f"UBICACIÓN ACTUAL DE GUSTAVO: coordenadas {lat:.5f}, {lng:.5f}"
            info_contexto += f"\n{info_ubicacion}"

            # Si pide clima, obtenerlo para su ubicación real
            if not clima:
                info_contexto += f"\n(Usar coordenadas {lat},{lng} para clima si se pregunta)"
    hechos = cargar_hechos()
    conversaciones = cargar_conversaciones(10)

    try:
        conn = get_conn()
        rows = conn.run("SELECT nombre, acciones FROM comandos ORDER BY id")
        conn.close()
        comandos = [{"nombre": r[0], "acciones": r[1]} for r in rows]
    except Exception:
        comandos = []

    system_prompt = PERSONALIDAD
    system_prompt += f"\n\n{info_contexto}"

    if hechos:
        system_prompt += "\n\nLo que sé de Gustavo:\n"
        for h in hechos:
            system_prompt += f"- {h}\n"

    if comandos:
        system_prompt += "\n\nComandos personalizados de Gustavo:\n"
        for c in comandos:
            system_prompt += f"- '{c['nombre']}': {c['acciones']}\n"
        system_prompt += "Cuando Gustavo diga el nombre de un comando, ejecuta sus acciones.\n"

    if conversaciones:
        system_prompt += "\n\nConversaciones recientes:\n"
        for conv in conversaciones:
            u = conv['usuario'][:200] if len(conv['usuario']) > 200 else conv['usuario']
            a = conv['akuzfiro'][:300] if len(conv['akuzfiro']) > 300 else conv['akuzfiro']
            system_prompt += f"Gustavo: {u}\nAkuzfiro: {a}\n"

    if not imagen_b64 and necesita_busqueda(mensaje):
        info_web = buscar_web(mensaje)
        if info_web:
            system_prompt += f"\n\nINFORMACIÓN WEB ENCONTRADA:\n{info_web}"

    # Construir el mensaje de usuario — con o sin imagen
    if imagen_b64:
        # Modelo de visión — mensaje multimodal
        user_content = [
            {
                "type": "text",
                "text": mensaje
            },
            {
                "type": "image_url",
                "image_url": {
                    "url": f"data:{imagen_mime};base64,{imagen_b64}"
                }
            }
        ]
        modelo = "meta-llama/llama-4-scout-17b-16e-instruct"
        tokens_max = 1024
    else:
        user_content = mensaje
        palabras_archivo = ["crea", "crear", "genera", "generar", "haz", "hacer", "excel", "pdf", "word", "powerpoint", "pptx", "documento", "presentacion", "presentación", "reporte", "bitacora", "bitácora"]
        palabras_largas = palabras_archivo + ["noticias", "noticia", "resumen de", "qué pasó", "últimas"]
        es_largo = any(p in mensaje.lower() for p in palabras_largas)
        tokens_max = 1500 if es_largo else 800
        modelo = "llama-3.3-70b-versatile"

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_content}
    ]

    try:
        response = client.chat.completions.create(
            model=modelo,
            messages=messages,
            temperature=0.85,
            max_tokens=tokens_max,
            frequency_penalty=0.3 if not imagen_b64 else 0,
            presence_penalty=0.1 if not imagen_b64 else 0
        )
        respuesta = response.choices[0].message.content

        # --- GUARDAR CALENDARIO EN MEMORIA si la imagen parece un calendario ---
        if imagen_b64:
            palabras_calendario = ["calendario", "agenda", "horario", "schedule", "fechas",
                                   "actividades", "eventos", "clases", "materias", "semestre"]
            es_calendario = any(p in mensaje.lower() for p in palabras_calendario) or \
                            any(p in respuesta.lower() for p in ["enero","febrero","marzo","abril","mayo",
                                "junio","julio","agosto","septiembre","octubre","noviembre","diciembre",
                                "lunes","martes","miércoles","jueves","viernes","semana","mes"])
            if es_calendario:
                try:
                    extractor = client.chat.completions.create(
                        model="meta-llama/llama-4-scout-17b-16e-instruct",
                        messages=[
                            {"role": "system", "content": (
                                "Extrae TODOS los eventos, fechas, actividades, clases o tareas de esta imagen de calendario. "
                                "Responde SOLO con una lista, un evento por línea, en formato: "
                                "'FECHA: descripción del evento'. "
                                "Si hay materias o clases con horario, incluye día y hora. "
                                "Sé específico y completo. Sin explicaciones adicionales."
                            )},
                            {"role": "user", "content": [
                                {"type": "text", "text": "Extrae todos los eventos de este calendario:"},
                                {"type": "image_url", "image_url": {"url": f"data:{imagen_mime};base64,{imagen_b64}"}}
                            ]}
                        ],
                        temperature=0.0,
                        max_tokens=1000
                    )
                    eventos_texto = extractor.choices[0].message.content.strip()
                    if eventos_texto and len(eventos_texto) > 10:
                        # Guardar cada evento como hecho
                        for linea in eventos_texto.split("\n"):
                            linea = linea.strip().lstrip("-•*·123456789. ")
                            if len(linea) > 5:
                                guardar_hecho(f"[CALENDARIO] {linea}")
                        # También guardar un resumen completo
                        guardar_hecho(f"[CALENDARIO COMPLETO] {eventos_texto[:800]}")
                        print(f"Calendario guardado: {len(eventos_texto.split(chr(10)))} eventos")
                except Exception as ex:
                    print(f"Error extrayendo calendario: {ex}")
        # --- FIN GUARDAR CALENDARIO ---

        # --- POST-PROCESO: inyectar tags [LUGAR] si el modelo no los incluyó ---
        if "[LUGAR]" not in respuesta and not imagen_b64:
            # No inyectar GPS si la respuesta es sobre un gasto/compra/producto
            palabras_gasto = ["anotado", "compraste", "gastaste", "pagaste", "compra",
                              "gasto", "precio", "pesos", "moneda", "costo", "vale",
                              "empeño", "préstamo", "prestamo", "ganancia", "venta"]
            es_respuesta_gasto = any(p in respuesta.lower() for p in palabras_gasto)

            palabras_lugar = ["recomiendo", "recomiend", "lugar", "restaurante", "taquería", "taqueria",
                              "café", "cafe", "parque", "hospital", "hotel", "bar", "cantina",
                              "plaza", "mercado", "farmacia", "clínica", "clinica", "gym", "gimnasio",
                              "cine", "teatro", "museo", "heladería", "heladeria", "panadería", "panaderia",
                              "te sugiero", "puedes ir", "puedes visitar", "visita", "conoce", "también está",
                              "también puedes", "otro lugar", "un lugar"]
            if not es_respuesta_gasto and any(p in respuesta.lower() for p in palabras_lugar):
                # Extraer nombres entre comillas — ej: "Tacos El Güero", 'La Michoacana'
                nombres_entre_comillas = re.findall(r'["\u201c\u201d\u2018\u2019\u00ab\u00bb]([^"\u201c\u201d\u2018\u2019\u00ab\u00bb]{3,60})["\u201c\u201d\u2018\u2019\u00ab\u00bb]', respuesta)

                # Filtrar solo los que parecen nombres de lugares (capitalizados, no frases genéricas)
                ignorar = {"xalapa", "veracruz", "méxico", "mexico", "la ciudad", "este lugar",
                           "estos lugares", "un lugar", "el lugar", "la zona", "el área"}
                lugares_encontrados = []
                vistos = set()
                for nombre in nombres_entre_comillas:
                    nombre = nombre.strip()
                    if nombre.lower() in ignorar or len(nombre) < 3:
                        continue
                    # Debe tener al menos una palabra capitalizada
                    if any(w[0].isupper() for w in nombre.split() if w):
                        key = nombre.lower()
                        if key not in vistos:
                            vistos.add(key)
                            lugares_encontrados.append(nombre)

                if lugares_encontrados:
                    # Determinar ciudad — buscar ciudad mencionada en el texto
                    ciudades_conocidas = ["xalapa", "veracruz", "cdmx", "ciudad de méxico",
                                         "guadalajara", "monterrey", "puebla", "oaxaca"]
                    ciudad_detectada = "Xalapa, Veracruz, México"
                    for c in ciudades_conocidas:
                        if c in respuesta.lower():
                            if c == "xalapa":
                                ciudad_detectada = "Xalapa, Veracruz, México"
                            elif c in ("veracruz",):
                                ciudad_detectada = "Veracruz, México"
                            elif c in ("cdmx", "ciudad de méxico"):
                                ciudad_detectada = "Ciudad de México, México"
                            break

                    tags = ""
                    for nombre in lugares_encontrados:
                        tags += f'[LUGAR]{json.dumps({"nombre": nombre, "ciudad": ciudad_detectada}, ensure_ascii=False)}[/LUGAR]\n'

                    # Quitar pregunta molesta de dirección/ubicación
                    respuesta = re.sub(
                        r'\s*[¿]?(quieres|deseas|te\s+gustar[ií]a)\s+(saber\s+)?(que\s+)?(te\s+)?(d[eéi][^\n.?]{0,80}|saber\s+)?(direcci[oó]n|ubicaci[oó]n|ubicad[oa]s?|c[oó]mo\s+llegar|indicaci[oó]n)[^\n.?]{0,40}[?]?',
                        "",
                        respuesta,
                        flags=re.IGNORECASE
                    ).strip()
                    respuesta = respuesta.rstrip("?¿ \n") + "\n" + tags
        # --- FIN POST-PROCESO ---

        guardar_conversacion(mensaje, respuesta)

        try:
            conn = get_conn()
            count = conn.run("SELECT COUNT(*) FROM conversaciones")[0][0]
            conn.close()
            if count % 3 == 0:
                extraer_hechos_automatico(mensaje, respuesta)
        except Exception:
            pass

        audio_b64_resp = None
        if con_voz and ELEVENLABS_API_KEY:
            audio_bytes = texto_a_voz(respuesta)
            if audio_bytes:
                audio_b64_resp = base64.b64encode(audio_bytes).decode("utf-8")

        return jsonify({"respuesta": respuesta, "audio": audio_b64_resp})

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/memoria", methods=["GET"])
def ver_memoria():
    return jsonify({
        "hechos": cargar_hechos(),
        "conversaciones": cargar_conversaciones(50)
    })

@app.route("/limpiar-hechos", methods=["POST"])
def limpiar_hechos():
    try:
        conn = get_conn()
        conn.run("DELETE FROM hechos WHERE hecho LIKE '%NINGUNO%' OR hecho LIKE '%no se menciona%' OR hecho LIKE '%no hay%' OR hecho LIKE '%la respuesta%' OR LENGTH(hecho) < 6")
        conn.close()
        return jsonify({"ok": True, "mensaje": "Hechos basura eliminados"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/hecho", methods=["POST"])
def agregar_hecho():
    data = request.json
    hecho = data.get("hecho", "")
    if hecho:
        guardar_hecho(hecho)
        return jsonify({"ok": True})
    return jsonify({"error": "Hecho vacío"}), 400

@app.route("/tts", methods=["POST"])
def tts():
    """Convierte texto a voz con gTTS y devuelve MP3."""
    try:
        data = request.json
        texto = data.get("texto", "")
        if not texto:
            return jsonify({"error": "Texto vacío"}), 400
        # Limpiar URLs y markdown
        texto_limpio = re.sub(r'https?://\S+', '', texto)
        texto_limpio = re.sub(r'\*\*|__|\*|_|`', '', texto_limpio)
        texto_limpio = texto_limpio.strip()[:800]
        buf = io.BytesIO()
        tts_obj = gTTS(text=texto_limpio, lang='es', tld='com.mx', slow=False)
        tts_obj.write_to_fp(buf)
        buf.seek(0)
        return send_file(buf, mimetype="audio/mpeg")
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/analizar-imagen", methods=["POST"])
def analizar_imagen():
    """
    Recibe una imagen en base64 y una pregunta opcional.
    Usa llama-4-scout (visión) para analizarla.
    """
    try:
        data = request.json
        imagen_b64 = data.get("imagen", "")      # data:image/jpeg;base64,....
        pregunta   = data.get("pregunta", "").strip()
        if not imagen_b64:
            return jsonify({"error": "Sin imagen"}), 400

        # Asegurar que viene con el prefijo data URI
        if not imagen_b64.startswith("data:"):
            imagen_b64 = "data:image/jpeg;base64," + imagen_b64

        if not pregunta:
            pregunta = ("Analiza esta imagen en detalle. Describe qué ves, "
                        "identifica objetos, texto, personas, animales, plantas "
                        "o cualquier cosa relevante. Si hay texto, léelo completo. "
                        "Responde en español de forma natural y directa.")

        tz_mexico = pytz.timezone("America/Mexico_City")
        ahora = datetime.now(tz_mexico).strftime("%A %d de %B de %Y, %H:%M hrs")

        system_vision = (PERSONALIDAD.split("REGLA #1")[0].strip() +
                         f"\n\nFECHA Y HORA: {ahora} (Xalapa, Veracruz, México)")

        messages = [
            {"role": "system", "content": system_vision},
            {
                "role": "user",
                "content": [
                    {"type": "text",      "text": pregunta},
                    {"type": "image_url", "image_url": {"url": imagen_b64}}
                ]
            }
        ]

        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=messages,
            temperature=0.7,
            max_tokens=1200
        )
        respuesta = response.choices[0].message.content
        guardar_conversacion(f"[IMAGEN] {pregunta}", respuesta)
        return jsonify({"respuesta": respuesta})

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/favicon.ico")
def favicon():
    """Evita el error 404 del favicon."""
    return "", 204

@app.route("/icon-192.png")
@app.route("/icon-512.png")
def icon():
    """Genera ícono SVG convertido a PNG para la PWA."""
    size = 512 if "512" in request.path else 192
    # Generar PNG simple con reportlab
    buf = io.BytesIO()
    from reportlab.graphics import renderPM
    from reportlab.graphics.shapes import Drawing, Circle, String
    from reportlab.lib import colors as rl_colors

    d = Drawing(size, size)
    # Fondo
    bg = Circle(size//2, size//2, size//2)
    bg.fillColor = rl_colors.HexColor("#050508")
    bg.strokeColor = None
    d.add(bg)
    # Círculo acento
    ring = Circle(size//2, size//2, size//2 - size//10)
    ring.fillColor = None
    ring.strokeColor = rl_colors.HexColor("#00d4ff")
    ring.strokeWidth = size//20
    d.add(ring)
    # Letra A
    font_size = size // 2
    txt = String(size//2, size//4, "A",
                 fontSize=font_size,
                 fillColor=rl_colors.HexColor("#00d4ff"),
                 textAnchor="middle")
    d.add(txt)

    renderPM.drawToFile(d, buf, fmt="PNG")
    buf.seek(0)
    return send_file(buf, mimetype="image/png")


@app.route("/comandos", methods=["GET"])
def ver_comandos():
    try:
        conn = get_conn()
        rows = conn.run("SELECT nombre, acciones FROM comandos ORDER BY id")
        conn.close()
        return jsonify([{"nombre": r[0], "acciones": r[1]} for r in rows])
    except Exception:
        return jsonify([])

@app.route("/comandos", methods=["POST"])
def guardar_comando():
    try:
        data = request.json
        nombre = data.get("nombre", "").strip()
        acciones = data.get("acciones", "").strip()
        if not nombre or not acciones:
            return jsonify({"error": "Faltan datos"}), 400
        conn = get_conn()
        conn.run("INSERT INTO comandos (nombre, acciones) VALUES (:n, :a) ON CONFLICT (nombre) DO UPDATE SET acciones = :a",
                 n=nombre, a=acciones)
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/comandos/<nombre>", methods=["DELETE"])
def eliminar_comando(nombre):
    try:
        conn = get_conn()
        conn.run("DELETE FROM comandos WHERE nombre = :n", n=nombre)
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/gastos", methods=["GET"])
def ver_gastos():
    try:
        conn = get_conn()
        rows = conn.run("""
            SELECT id, fecha, descripcion, monto, categoria
            FROM gastos ORDER BY fecha DESC LIMIT 100
        """)
        conn.close()
        total = sum(float(r[3]) for r in rows)
        gastos = [{"id": r[0], "fecha": str(r[1])[:16], "descripcion": r[2],
                   "monto": float(r[3]), "categoria": r[4]} for r in rows]
        return jsonify({"gastos": gastos, "total": total})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/gastos", methods=["POST"])
def agregar_gasto():
    try:
        data = request.json
        descripcion = data.get("descripcion", "").strip()
        monto = float(data.get("monto", 0))
        categoria = data.get("categoria", "general").strip()
        if not descripcion or monto <= 0:
            return jsonify({"error": "Datos inválidos"}), 400
        conn = get_conn()
        conn.run(
            "INSERT INTO gastos (descripcion, monto, categoria) VALUES (:d, :m, :c)",
            d=descripcion, m=monto, c=categoria
        )
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/gastos/<int:gasto_id>", methods=["DELETE"])
def eliminar_gasto(gasto_id):
    try:
        conn = get_conn()
        conn.run("DELETE FROM gastos WHERE id = :id", id=gasto_id)
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# --- GANANCIAS ---
@app.route("/ganancias", methods=["POST"])
def agregar_ganancia():
    try:
        data = request.json
        descripcion = data.get("descripcion", "").strip()
        monto = float(data.get("monto", 0))
        categoria = data.get("categoria", "otro").strip()
        if not descripcion or monto <= 0:
            return jsonify({"error": "Datos inválidos"}), 400
        conn = get_conn()
        conn.run("INSERT INTO ganancias (descripcion, monto, categoria) VALUES (:d, :m, :c)",
                 d=descripcion, m=monto, c=categoria)
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/ganancias", methods=["GET"])
def ver_ganancias():
    try:
        conn = get_conn()
        rows = conn.run("SELECT id, fecha, descripcion, monto, categoria FROM ganancias ORDER BY fecha DESC LIMIT 50")
        conn.close()
        return jsonify([{"id": r[0], "fecha": str(r[1])[:16], "descripcion": r[2], "monto": float(r[3]), "categoria": r[4]} for r in rows])
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# --- VENTAS ---
@app.route("/ventas", methods=["POST"])
def agregar_venta():
    try:
        data = request.json
        descripcion = data.get("descripcion", "").strip()
        monto = float(data.get("monto", 0))
        if not descripcion or monto <= 0:
            return jsonify({"error": "Datos inválidos"}), 400
        conn = get_conn()
        conn.run("INSERT INTO ventas (descripcion, monto) VALUES (:d, :m)", d=descripcion, m=monto)
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# --- PRÉSTAMOS ---
@app.route("/prestamos", methods=["POST"])
def agregar_prestamo():
    try:
        data = request.json
        tipo = data.get("tipo", "dado")  # "dado" o "recibido"
        persona = data.get("persona", "").strip()
        monto = float(data.get("monto", 0))
        descripcion = data.get("descripcion", "").strip()
        if not persona or monto <= 0:
            return jsonify({"error": "Datos inválidos"}), 400
        conn = get_conn()
        conn.run("INSERT INTO prestamos (tipo, persona, monto, descripcion) VALUES (:t, :p, :m, :d)",
                 t=tipo, p=persona, m=monto, d=descripcion)
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/prestamos", methods=["GET"])
def ver_prestamos():
    try:
        conn = get_conn()
        rows = conn.run("SELECT id, fecha, tipo, persona, monto, descripcion, saldado FROM prestamos ORDER BY fecha DESC LIMIT 50")
        conn.close()
        return jsonify([{"id": r[0], "fecha": str(r[1])[:16], "tipo": r[2], "persona": r[3], "monto": float(r[4]), "descripcion": r[5], "saldado": r[6]} for r in rows])
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# --- EMPEÑOS ---
@app.route("/empenos", methods=["POST"])
def agregar_empeno():
    try:
        data = request.json
        descripcion = data.get("descripcion", "").strip()
        monto = float(data.get("monto", 0))
        lugar = data.get("lugar", "").strip()
        if not descripcion:
            return jsonify({"error": "Datos inválidos"}), 400
        conn = get_conn()
        conn.run("INSERT INTO empenos (descripcion, monto, lugar) VALUES (:d, :m, :l)",
                 d=descripcion, m=monto, l=lugar)
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# --- LISTA DE COMPRAS ---
@app.route("/lista-compras", methods=["POST"])
def agregar_lista_compras():
    try:
        data = request.json
        items = data.get("items", [])
        if not items:
            return jsonify({"error": "Lista vacía"}), 400
        conn = get_conn()
        for item in items:
            if item.strip():
                conn.run("INSERT INTO lista_compras (item) VALUES (:i)", i=item.strip())
        conn.close()
        return jsonify({"ok": True, "agregados": len(items)})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/lista-compras", methods=["GET"])
def ver_lista_compras():
    try:
        conn = get_conn()
        rows = conn.run("SELECT id, item, comprado, fecha FROM lista_compras WHERE comprado = FALSE ORDER BY fecha DESC")
        conn.close()
        return jsonify([{"id": r[0], "item": r[1], "comprado": r[2]} for r in rows])
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/lista-compras/<int:item_id>/comprado", methods=["POST"])
def marcar_comprado(item_id):
    try:
        conn = get_conn()
        conn.run("UPDATE lista_compras SET comprado = TRUE WHERE id = :id", id=item_id)
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# --- BITÁCORA ---
@app.route("/bitacora", methods=["POST"])
def agregar_bitacora():
    try:
        data = request.json
        descripcion = data.get("descripcion", "").strip()
        categoria = data.get("categoria", "general").strip()
        intervalo_dias = data.get("intervalo_dias", None)
        if not descripcion:
            return jsonify({"error": "Descripción vacía"}), 400
        tz_mexico = pytz.timezone("America/Mexico_City")
        ahora = datetime.now(tz_mexico)
        proxima = None
        if intervalo_dias:
            from datetime import timedelta
            proxima = ahora + timedelta(days=int(intervalo_dias))
        conn = get_conn()
        conn.run(
            "INSERT INTO bitacora (descripcion, categoria, proxima_fecha, intervalo_dias) VALUES (:d, :c, :p, :i)",
            d=descripcion, c=categoria,
            p=proxima.replace(tzinfo=None) if proxima else None,
            i=int(intervalo_dias) if intervalo_dias else None
        )
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/bitacora", methods=["GET"])
def ver_bitacora():
    try:
        conn = get_conn()
        rows = conn.run("""
            SELECT id, fecha, descripcion, categoria, proxima_fecha, intervalo_dias
            FROM bitacora ORDER BY fecha DESC LIMIT 100
        """)
        conn.close()
        result = []
        for r in rows:
            result.append({
                "id": r[0], "fecha": str(r[1])[:16],
                "descripcion": r[2], "categoria": r[3],
                "proxima_fecha": str(r[4])[:16] if r[4] else None,
                "intervalo_dias": r[5]
            })
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/bitacora/proximas", methods=["GET"])
def bitacora_proximas():
    """Entradas de bitácora cuya próxima fecha ya llegó o está próxima (3 días)."""
    try:
        from datetime import timedelta
        tz_mexico = pytz.timezone("America/Mexico_City")
        ahora = datetime.now(tz_mexico).replace(tzinfo=None)
        pronto = ahora + timedelta(days=3)
        conn = get_conn()
        rows = conn.run("""
            SELECT id, descripcion, categoria, proxima_fecha
            FROM bitacora
            WHERE proxima_fecha IS NOT NULL AND proxima_fecha <= :pronto
            ORDER BY proxima_fecha ASC
        """, pronto=pronto)
        conn.close()
        return jsonify([{
            "id": r[0], "descripcion": r[1],
            "categoria": r[2], "proxima_fecha": str(r[3])[:16]
        } for r in rows])
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# --- RESUMEN FINANCIERO ---
@app.route("/finanzas/resumen", methods=["GET"])
def resumen_finanzas():
    try:
        conn = get_conn()
        gastos = conn.run("SELECT COALESCE(SUM(monto),0) FROM gastos WHERE fecha >= date_trunc('month', NOW())")[0][0]
        ganancias = conn.run("SELECT COALESCE(SUM(monto),0) FROM ganancias WHERE fecha >= date_trunc('month', NOW())")[0][0]
        ventas = conn.run("SELECT COALESCE(SUM(monto),0) FROM ventas WHERE fecha >= date_trunc('month', NOW())")[0][0]
        prestamos_dados = conn.run("SELECT COALESCE(SUM(monto),0) FROM prestamos WHERE tipo='dado' AND saldado=FALSE")[0][0]
        prestamos_recibidos = conn.run("SELECT COALESCE(SUM(monto),0) FROM prestamos WHERE tipo='recibido' AND saldado=FALSE")[0][0]
        empenos = conn.run("SELECT COALESCE(SUM(monto),0) FROM empenos WHERE recuperado=FALSE")[0][0]
        compras_pendientes = conn.run("SELECT COUNT(*) FROM lista_compras WHERE comprado=FALSE")[0][0]
        conn.close()
        return jsonify({
            "gastos_mes": float(gastos),
            "ganancias_mes": float(ganancias),
            "ventas_mes": float(ventas),
            "balance_mes": float(ganancias) + float(ventas) - float(gastos),
            "prestamos_dados": float(prestamos_dados),
            "prestamos_recibidos": float(prestamos_recibidos),
            "empenos_activos": float(empenos),
            "compras_pendientes": int(compras_pendientes)
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/recordatorios", methods=["GET"])
def ver_recordatorios():
    try:
        tz_mexico = pytz.timezone("America/Mexico_City")
        conn = get_conn()
        rows = conn.run("""
            SELECT id, hora_aviso, mensaje, completado
            FROM recordatorios WHERE completado = FALSE
            ORDER BY hora_aviso ASC
        """)
        conn.close()
        return jsonify([{
            "id": r[0],
            "hora_aviso": str(r[1])[:16],
            "mensaje": r[2],
            "completado": r[3]
        } for r in rows])
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/recordatorios", methods=["POST"])
def agregar_recordatorio():
    try:
        data = request.json
        mensaje = data.get("mensaje", "").strip()
        hora_aviso = data.get("hora_aviso", "")
        if not mensaje or not hora_aviso:
            return jsonify({"error": "Faltan datos"}), 400
        conn = get_conn()
        conn.run(
            "INSERT INTO recordatorios (hora_aviso, mensaje) VALUES (:h, :m)",
            h=hora_aviso, m=mensaje
        )
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/recordatorios/<int:rec_id>/completar", methods=["POST"])
def completar_recordatorio(rec_id):
    try:
        conn = get_conn()
        conn.run("UPDATE recordatorios SET completado = TRUE WHERE id = :id", id=rec_id)
        conn.close()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/recordatorios/pendientes", methods=["GET"])
def recordatorios_pendientes():
    """Verifica si hay recordatorios que deben sonar ahora (ventana de 5 minutos)."""
    try:
        from datetime import timedelta
        tz_mexico = pytz.timezone("America/Mexico_City")
        ahora = datetime.now(tz_mexico)
        hace5min = ahora - timedelta(minutes=5)
        conn = get_conn()
        rows = conn.run("""
            SELECT id, mensaje FROM recordatorios
            WHERE completado = FALSE
            AND hora_aviso <= :ahora
            AND hora_aviso >= :hace5min
        """, ahora=ahora, hace5min=hace5min)
        conn.close()
        return jsonify([{"id": r[0], "mensaje": r[1]} for r in rows])
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/recordatorios/parsear", methods=["POST"])
def parsear_recordatorio():
    """
    Recibe una frase natural (ej: 'avísame a las 3pm que tengo reunión')
    y devuelve hora_aviso en formato ISO y el mensaje extraído.
    """
    try:
        from datetime import timedelta
        data = request.json
        frase = data.get("frase", "").strip()
        if not frase:
            return jsonify({"error": "Frase vacía"}), 400

        tz_mexico = pytz.timezone("America/Mexico_City")
        ahora = datetime.now(tz_mexico)

        prompt = f"""Extrae la hora y el mensaje de este recordatorio. La fecha/hora actual es: {ahora.strftime('%Y-%m-%d %H:%M')} (Xalapa, México, hora del centro).

Frase: "{frase}"

Responde SOLO con JSON válido en este formato exacto, sin explicaciones:
{{"hora_aviso": "YYYY-MM-DD HH:MM:SS", "mensaje": "texto del recordatorio"}}

Reglas:
- Si dice "en X minutos", suma esos minutos al tiempo actual
- Si dice "a las 3pm" o "a las 15:00", usa esa hora de HOY (si ya pasó, usa mañana)
- Si dice "mañana a las...", usa la fecha de mañana
- Si no especifica AM/PM y la hora es < 8, asume PM
- El mensaje debe ser corto y claro, sin la parte de "avísame" o "recuérdame"
- Solo JSON, nada más"""

        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.0,
            max_tokens=100
        )
        resultado = response.choices[0].message.content.strip()
        # Limpiar si viene con backticks
        resultado = resultado.replace("```json", "").replace("```", "").strip()
        parsed = json.loads(resultado)
        return jsonify(parsed)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# Cache temporal para descargas (en memoria, se limpia al reiniciar)
_download_cache = {}

@app.route("/preparar-descarga", methods=["POST"])
def preparar_descarga():
    """Guarda el payload en cache y devuelve un token para descarga GET."""
    import uuid
    data = request.json
    tipo = data.get("tipo", "word")
    payload = data.get("payload", {})
    token = str(uuid.uuid4())[:8]
    _download_cache[token] = {"tipo": tipo, "payload": payload}
    return jsonify({"token": token, "url": f"/descargar/{token}"})

@app.route("/descargar/<token>", methods=["GET"])
def descargar_archivo(token):
    """Descarga el archivo usando el token generado."""
    if token not in _download_cache:
        return jsonify({"error": "Token inválido o expirado"}), 404
    item = _download_cache.pop(token)
    tipo = item["tipo"]
    payload = item["payload"]
    # Redirigir al generador correspondiente
    from flask import make_response
    if tipo == "word":
        with app.test_request_context(json=payload):
            resp = generar_word()
            return resp
    elif tipo == "excel":
        with app.test_request_context(json=payload):
            return generar_excel()
    elif tipo == "pdf":
        with app.test_request_context(json=payload):
            return generar_pdf()
    elif tipo == "pptx":
        with app.test_request_context(json=payload):
            return generar_pptx()
    return jsonify({"error": "Tipo no soportado"}), 400


@app.route("/generar-word", methods=["POST"])
def generar_word():
    try:
        data = request.json
        titulo = data.get("titulo", "Documento")
        contenido = data.get("contenido", "")
        secciones = data.get("secciones", [])

        doc = Document()

        # Estilos
        estilo_normal = doc.styles["Normal"]
        estilo_normal.font.name = "Calibri"
        estilo_normal.font.size = Pt(11)

        # Título
        titulo_par = doc.add_heading(titulo, level=0)
        titulo_par.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in titulo_par.runs:
            run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
            run.font.size = Pt(18)

        doc.add_paragraph()

        # Contenido principal
        if contenido:
            for linea in contenido.split("\n"):
                if linea.strip():
                    p = doc.add_paragraph(linea.strip())
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        # Secciones
        for seccion in secciones:
            doc.add_paragraph()
            h = doc.add_heading(seccion.get("titulo", ""), level=1)
            for run in h.runs:
                run.font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)
            texto_sec = seccion.get("contenido", "")
            for linea in texto_sec.split("\n"):
                if linea.strip():
                    p = doc.add_paragraph(linea.strip())
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

            # Tabla de sección si existe
            tabla_sec = seccion.get("tabla", None)
            if tabla_sec:
                encabezados = tabla_sec.get("encabezados", [])
                filas = tabla_sec.get("filas", [])
                if encabezados:
                    t = doc.add_table(rows=1, cols=len(encabezados))
                    t.style = "Table Grid"
                    hdr = t.rows[0].cells
                    for i, enc in enumerate(encabezados):
                        hdr[i].text = enc
                        hdr[i].paragraphs[0].runs[0].font.bold = True
                        hdr[i].paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    for fila in filas:
                        row = t.add_row().cells
                        for i, val in enumerate(fila):
                            if i < len(row):
                                row[i].text = str(val)

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        nombre_archivo = f"{titulo.replace(' ', '_')}.docx"
        return send_file(buf,
                        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        as_attachment=True, download_name=nombre_archivo)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/generar-pptx", methods=["POST"])
def generar_pptx():
    try:
        data = request.json
        titulo = data.get("titulo", "Presentación")
        diapositivas = data.get("diapositivas", [])

        prs = Presentation()
        prs.slide_width = Emu(9144000)
        prs.slide_height = Emu(5143500)

        COLOR_FONDO = PRGBColor(0x1F, 0x38, 0x64)
        COLOR_ACENTO = PRGBColor(0x2E, 0x75, 0xB6)
        COLOR_TEXTO = PRGBColor(0xFF, 0xFF, 0xFF)
        COLOR_SUBTEXTO = PRGBColor(0xD6, 0xE4, 0xF0)

        def set_bg(slide, color):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = color

        # Diapositiva de título
        lay_blank = prs.slide_layouts[6]
        slide_titulo = prs.slides.add_slide(lay_blank)
        set_bg(slide_titulo, COLOR_FONDO)

        txb = slide_titulo.shapes.add_textbox(PInches(0.5), PInches(1.5), PInches(9), PInches(1.5))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = titulo
        run.font.size = PPt(40)
        run.font.bold = True
        run.font.color.rgb = COLOR_TEXTO

        # Línea decorativa
        line = slide_titulo.shapes.add_shape(1, PInches(2), PInches(3.2), PInches(5), Emu(50000))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACENTO
        line.line.fill.background()

        # Diapositivas de contenido
        for i, diap in enumerate(diapositivas):
            slide = prs.slides.add_slide(lay_blank)
            set_bg(slide, COLOR_FONDO)

            # Número de diapositiva
            num_txb = slide.shapes.add_textbox(PInches(8.5), PInches(0.1), PInches(0.5), PInches(0.3))
            num_tf = num_txb.text_frame
            num_p = num_tf.paragraphs[0]
            num_run = num_p.add_run()
            num_run.text = str(i + 1)
            num_run.font.size = PPt(12)
            num_run.font.color.rgb = COLOR_SUBTEXTO

            # Título de diapositiva
            titulo_diap = diap.get("titulo", "")
            txb_t = slide.shapes.add_textbox(PInches(0.4), PInches(0.3), PInches(8.5), PInches(0.8))
            tf_t = txb_t.text_frame
            p_t = tf_t.paragraphs[0]
            run_t = p_t.add_run()
            run_t.text = titulo_diap
            run_t.font.size = PPt(28)
            run_t.font.bold = True
            run_t.font.color.rgb = PRGBColor(0x00, 0xD4, 0xFF)

            # Línea bajo título
            sep = slide.shapes.add_shape(1, PInches(0.4), PInches(1.1), PInches(8.5), Emu(40000))
            sep.fill.solid()
            sep.fill.fore_color.rgb = COLOR_ACENTO
            sep.line.fill.background()

            # Contenido
            puntos = diap.get("puntos", [])
            contenido_texto = diap.get("contenido", "")

            txb_c = slide.shapes.add_textbox(PInches(0.4), PInches(1.3), PInches(8.5), PInches(3.5))
            tf_c = txb_c.text_frame
            tf_c.word_wrap = True

            if puntos:
                for j, punto in enumerate(puntos):
                    p_c = tf_c.paragraphs[0] if j == 0 else tf_c.add_paragraph()
                    p_c.space_before = PPt(6)
                    run_c = p_c.add_run()
                    run_c.text = f"• {punto}"
                    run_c.font.size = PPt(18)
                    run_c.font.color.rgb = COLOR_SUBTEXTO
            elif contenido_texto:
                p_c = tf_c.paragraphs[0]
                run_c = p_c.add_run()
                run_c.text = contenido_texto
                run_c.font.size = PPt(18)
                run_c.font.color.rgb = COLOR_SUBTEXTO

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        nombre_archivo = f"{titulo.replace(' ', '_')}.pptx"
        return send_file(buf,
                        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        as_attachment=True, download_name=nombre_archivo)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/generar-excel", methods=["POST"])
def generar_excel():
    try:
        data = request.json
        titulo = data.get("titulo", "Documento")
        encabezados = data.get("encabezados", [])
        filas = data.get("filas", [])
        secciones = data.get("secciones", [])  # Para documentos con secciones/grupos

        wb = Workbook()
        ws = wb.active
        ws.title = titulo[:31]

        # Estilos
        estilo_titulo = Font(name="Calibri", size=14, bold=True, color="FFFFFF")
        fill_titulo = PatternFill("solid", fgColor="1F3864")
        estilo_header = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        fill_header = PatternFill("solid", fgColor="2E75B6")
        fill_seccion = PatternFill("solid", fgColor="D6E4F0")
        estilo_seccion = Font(name="Calibri", size=11, bold=True, color="1F3864")
        borde = Border(
            left=Side(style="thin"), right=Side(style="thin"),
            top=Side(style="thin"), bottom=Side(style="thin")
        )
        alineacion_centro = Alignment(horizontal="center", vertical="center", wrap_text=True)

        fila_actual = 1

        # Título principal
        num_cols = max(len(encabezados), 1)
        ws.merge_cells(f"A{fila_actual}:{get_column_letter(num_cols)}{fila_actual}")
        celda_titulo = ws.cell(row=fila_actual, column=1, value=titulo)
        celda_titulo.font = estilo_titulo
        celda_titulo.fill = fill_titulo
        celda_titulo.alignment = alineacion_centro
        ws.row_dimensions[fila_actual].height = 30
        fila_actual += 1

        # Encabezados
        if encabezados:
            for col, enc in enumerate(encabezados, 1):
                c = ws.cell(row=fila_actual, column=col, value=enc)
                c.font = estilo_header
                c.fill = fill_header
                c.alignment = alineacion_centro
                c.border = borde
            ws.row_dimensions[fila_actual].height = 20
            fila_actual += 1

        # Filas simples
        for i, fila in enumerate(filas):
            fill_fila = PatternFill("solid", fgColor="EBF3FB" if i % 2 == 0 else "FFFFFF")
            for col, val in enumerate(fila, 1):
                c = ws.cell(row=fila_actual, column=col, value=val)
                c.fill = fill_fila
                c.border = borde
                c.alignment = Alignment(vertical="center", wrap_text=True)
            fila_actual += 1

        # Secciones (grupos con título y filas)
        for seccion in secciones:
            nombre_sec = seccion.get("nombre", "")
            filas_sec = seccion.get("filas", [])

            # Título de sección
            ws.merge_cells(f"A{fila_actual}:{get_column_letter(num_cols)}{fila_actual}")
            c = ws.cell(row=fila_actual, column=1, value=nombre_sec)
            c.font = estilo_seccion
            c.fill = fill_seccion
            c.alignment = alineacion_centro
            c.border = borde
            ws.row_dimensions[fila_actual].height = 18
            fila_actual += 1

            for i, fila in enumerate(filas_sec):
                fill_fila = PatternFill("solid", fgColor="EBF3FB" if i % 2 == 0 else "FFFFFF")
                for col, val in enumerate(fila, 1):
                    c = ws.cell(row=fila_actual, column=col, value=val)
                    c.fill = fill_fila
                    c.border = borde
                    c.alignment = Alignment(vertical="center", wrap_text=True)
                fila_actual += 1

        # Ajustar ancho de columnas
        for col in range(1, num_cols + 1):
            max_len = 0
            for row in ws.iter_rows(min_col=col, max_col=col):
                for cell in row:
                    if cell.value:
                        max_len = max(max_len, len(str(cell.value)))
            ws.column_dimensions[get_column_letter(col)].width = min(max(max_len + 2, 12), 40)

        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)
        nombre_archivo = f"{titulo.replace(' ', '_')}.xlsx"
        return send_file(buf, mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        as_attachment=True, download_name=nombre_archivo)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/generar-pdf", methods=["POST"])
def generar_pdf():
    try:
        data = request.json
        titulo = data.get("titulo", "Documento")
        contenido = data.get("contenido", "")
        secciones = data.get("secciones", [])
        tabla = data.get("tabla", None)

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4,
                               rightMargin=2*cm, leftMargin=2*cm,
                               topMargin=2*cm, bottomMargin=2*cm)

        estilos = getSampleStyleSheet()
        estilo_titulo = ParagraphStyle("titulo", parent=estilos["Title"],
                                      fontSize=18, textColor=colors.HexColor("#1F3864"),
                                      spaceAfter=12, alignment=TA_CENTER)
        estilo_subtitulo = ParagraphStyle("subtitulo", parent=estilos["Heading2"],
                                         fontSize=13, textColor=colors.HexColor("#2E75B6"),
                                         spaceAfter=6)
        estilo_cuerpo = ParagraphStyle("cuerpo", parent=estilos["Normal"],
                                      fontSize=11, leading=16, spaceAfter=8,
                                      alignment=TA_JUSTIFY)

        elementos = []
        elementos.append(Paragraph(titulo, estilo_titulo))
        elementos.append(Spacer(1, 0.3*inch))

        if contenido:
            for parrafo in contenido.split("\n"):
                if parrafo.strip():
                    elementos.append(Paragraph(parrafo.strip(), estilo_cuerpo))

        for seccion in secciones:
            elementos.append(Spacer(1, 0.2*inch))
            elementos.append(Paragraph(seccion.get("titulo", ""), estilo_subtitulo))
            for p in seccion.get("contenido", "").split("\n"):
                if p.strip():
                    elementos.append(Paragraph(p.strip(), estilo_cuerpo))

        if tabla:
            elementos.append(Spacer(1, 0.2*inch))
            encabezados_tabla = tabla.get("encabezados", [])
            filas_tabla = tabla.get("filas", [])
            tabla_data = [encabezados_tabla] + filas_tabla
            t = Table(tabla_data, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2E75B6")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, 0), 11),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#EBF3FB"), colors.white]),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#AAAAAA")),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]))
            elementos.append(t)

        doc.build(elementos)
        buf.seek(0)
        nombre_archivo = f"{titulo.replace(' ', '_')}.pdf"
        return send_file(buf, mimetype="application/pdf",
                        as_attachment=True, download_name=nombre_archivo)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/")
def index():
    return app.send_static_file("index.html")

# Inicializar DB al arrancar
try:
    init_db()
    print("Base de datos lista.")
except Exception as e:
    print(f"Error iniciando DB: {e}")

if __name__ == "__main__":
    print("Akuzfiro iniciando...")
    app.run(debug=False, host="0.0.0.0", port=5000)
